#!/usr/bin/env python3
"""
Extract speaker notes from a PPTX and write to a plain text file.
"""
from pptx import Presentation
from pathlib import Path

P = Path(__file__).resolve().parent
PPTX = P / 'Part3_with_simulation_updated.pptx'
OUT = P / 'Part3_with_simulation_notes.txt'

if not PPTX.exists():
    raise SystemExit('PPTX not found: ' + str(PPTX))

prs = Presentation(str(PPTX))
with open(OUT, 'w', encoding='utf-8') as f:
    for i, slide in enumerate(prs.slides, start=1):
        f.write(f'Slide {i}: ')
        title = slide.shapes.title.text if slide.shapes.title else ''
        f.write(title + '\n')
        notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ''
        f.write('--- Notes ---\n')
        f.write(notes + '\n\n')

print('Wrote notes to', OUT)
