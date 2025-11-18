#!/usr/bin/env python3
"""
Create Part3_with_simulation.pptx with slides and speaker notes for Part 3.
This script uses python-pptx. Run: pip install python-pptx
"""
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception as e:
    raise SystemExit("python-pptx not installed. Install with: pip install python-pptx")

OUT_DIR = Path(__file__).resolve().parent
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_FILE = OUT_DIR / 'Part3_with_simulation.pptx'

# Speaker notes per slide (English), aligned with the revised script.
slides = [
    {
        'title': 'Preliminary Analysis & Findings — St. Martin Rd vs Chong San Rd',
        'notes': "Hello. In this section I will present our preliminary analysis comparing St. Martin Road and Chong San Road, and then describe how we will simulate two intervention scenarios — stop merging, and route splitting. For data, we used real-time ETA snapshots from public APIs, consolidated into CSV files. Our 1-hour experiment produced 120 snapshots and a consolidated CSV under `vibeCoding101/PartX_simulation/monitor_outputs_1hr/`. First, I’ll summarize key metrics and observed patterns.",
        'body': "Data overview and reference: `vibeCoding101/PartX_simulation/monitor_outputs_1hr/` (120 snapshots)."
    },
    {
        'title': 'Key metrics computed',
        'notes': "In our processing pipeline we compute several core metrics from each snapshot and ETA record. These include: wait time (wait_s = ETA - snapshot_ts), per-route ETA counts, the null-ETA rate, per-minute average wait, and a headway proxy derived from observed ETAs. We also compute a 'cleaned' average (mean_wait_cleaned) after trimming the top 1% of extreme waits to better represent typical conditions. These KPIs form the baseline we will compare against simulated scenarios.",
        'body': "KPIs: mean, median, 90th pct, null-rate, per-minute mean, route counts, mean_wait_cleaned"
    },
    {
        'title': 'Wait time distribution',
        'notes': "Here is the wait-time distribution for the monitored period (see `monitor_outputs_60min/analysis/wait_distribution.png` or `monitor_outputs_until_0830/analysis/wait_distribution.png`). The distribution is right-skewed: most ETAs cluster at short waits—roughly within one to three minutes—but there is a long tail of larger waits. These long-tail values suggest episodes of congestion, vehicle bunching, or delayed reporting by the API. To reduce the influence of these extremes on summary statistics, we also show a trimmed mean (mean_wait_cleaned) for a more representative central tendency.",
        'body': "Placeholder for: monitor_outputs_60min/analysis/wait_distribution.png"
    },
    {
        'title': 'Mean wait per minute (cleaned)',
        'notes': "This time-series plot (files: `mean_wait_cleaned.png` and `combined_mean_wait.png`) shows the minute-by-minute cleaned mean wait. You can also view 5- and 10-minute aggregated versions in the `counts_per_5min` and `counts_per_10min` folders to reduce noise. The main patterns are: average waits increase during peak windows; some snapshots show pronounced spikes indicating temporary instability; and aggregated views make these trends easier to interpret and compare across the two stops.",
        'body': "Placeholder for: mean_wait_cleaned.png, combined_mean_wait.png (counts_per_5min/*, counts_per_10min/*)"
    },
    {
        'title': 'Route counts (top 20)',
        'notes': "This bar chart (`route_counts_top20.png`) lists the most frequently observed routes in our monitoring window. A small set of routes appears far more frequently, indicating higher service frequency or route concentration near these stops. Importantly, if the same route appears in both stops within the same snapshot repeatedly, that can indicate short headways or bus bunching—an operational concern we consider when proposing stop consolidation or route adjustments.",
        'body': "Placeholder for: route_counts_top20.png"
    },
    {
        'title': 'Null ETA & data issues',
        'notes': "We observed a non-trivial number of records where `eta == null` in API responses; raw examples can be found in `vibeCoding101/PartX_simulation/monitor_outputs_1hr/` snapshot JSONs. Null ETAs prevent direct wait computation and reduce effective sample size. To handle this, we compute and report a null-rate KPI, exclude null entries from numerical wait calculations, and flag snapshots for further inspection. Addressing nulls could require longer sampling, operator collaboration, or alternative data sources.",
        'body': "Placeholder: example snapshot path: monitor_outputs_1hr/snapshot_*.json"
    },
    {
        'title': 'Simulation objectives & contract',
        'notes': "The high-level objective is to quantify how two interventions change our KPIs relative to baseline: (A) merging the two bus stops into a single stop, and (B) splitting a frequent route so that half the vehicles serve St. Martin and half serve Chong San. Contract / inputs: observed headways and ETA-derived arrival distributions from `monitor_outputs_*`, route frequency by route, and estimated dwell times (set from literature or sensitivity ranges). Outputs: simulated passenger wait distributions, per-route counts at each stop, 90th percentile wait, and null-rate proxy. We will run stochastic replications to estimate means and 95% confidence intervals for these KPIs.",
        'body': "Simulation inputs: monitor_outputs_* (headways, counts). Outputs: simulated wait distributions, KPI tables. Code reference: sim_kmb_stops.py"
    },
    {
        'title': 'Scenario A — Merge stops',
        'notes': "In Scenario A we model collapsing St. Martin and Chong San into one consolidated stop located at the midpoint. Operational assumptions: combined arrival streams for overlapping routes; slightly increased dwell time per stop due to higher boarding volume, modeled as dwell = base_dwell + alpha * arrivals. We will simulate multiple demand levels: current observed demand, +10% and +25% demand. Expected outputs to compare vs baseline: mean wait, 90th percentile wait, and passenger throughput. Hypothesis: merging reduces duplicated service but may increase dwell-driven delay; net effect depends on arrival regularity and boarding volume.",
        'body': "Placeholder: simulation diagram and parameter table (dwell model, demand levels)"
    },
    {
        'title': 'Scenario B — Split route (half/half)',
        'notes': "In Scenario B we reassign an identified high-frequency route so roughly half of its vehicles serve St. Martin and half serve Chong San. Implementation assumptions: split operates by vehicle assignment rather than passenger rerouting; headways for each sub-route increase roughly twofold for the split segment. We will vary splitting rules (strict alternating assignment vs probabilistic split) and test sensitivity to compliance. Key comparison metrics: mean wait at each stop, variance of wait, and frequency of same-route co-occurrence. Hypothesis: splitting reduces per-stop crowding and wait variance but may increase average wait if headways grow too large.",
        'body': "Placeholder: simulation variants and expected KPI outputs"
    },
    {
        'title': 'Model mechanics & validation',
        'notes': "Our simulator is a discrete-event model implemented in Python (see `sim_kmb_stops.py`). Core mechanics: event list of vehicle arrivals drawn from observed ETA-derived interarrival distributions; passenger arrivals modeled as Poisson with rate per-minute estimated from route counts; boarding consumes dwell time; service times and route assignments follow scenario rules. For validation, we compare simulator baseline output with empirical KPIs (mean wait, per-minute counts, route co-occurrence) to ensure the model reproduces observed behavior before running interventions.",
        'body': "Placeholder: validation plots (baseline empirical vs simulated)"
    },
    {
        'title': 'KPIs & visuals from simulation',
        'notes': "For each scenario we will present: baseline vs scenario KPI table with mean and 95% CI for mean wait and 90th percentile wait; overlaid wait-time histograms; time-series of mean wait per minute; and a small table reporting expected passenger throughput and dwell impact. Figures will be saved into `vibeCoding101/PartX_simulation/sim_outputs/<scenario>/` for reproducibility.",
        'body': "Placeholder: example KPI table and histogram templates"
    },
    {
        'title': 'Preliminary conclusions & next steps',
        'notes': "To summarize: empirical data shows long-tailed waits and signs of vehicle bunching—these motivate two interventions we will test with our simulator: merging stops and splitting routes. We will validate the simulator against observed KPIs, run sensitivity analyses for demand and dwell assumptions, and then compare KPI changes across scenarios. The next section will detail simulation results and policy recommendations. That concludes Part 3.",
        'body': "Closing slide: next steps and link to Part 4"
    }
]

prs = Presentation()
# Use a simple title slide layout for the first slide
for i, s in enumerate(slides):
    if i == 0:
        slide_layout = prs.slide_layouts[0]
    else:
        slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    # set title
    if slide.shapes.title:
        slide.shapes.title.text = s['title']
    # add body placeholder (textbox)
    left = Inches(0.5)
    top = Inches(1.6)
    width = Inches(9)
    height = Inches(3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = s['body']
    p.font.size = Pt(14)
    # add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = s['notes']

# remove the initial default empty first slide that Presentation() creates (if any)
# Actually Presentation() has an initial slide count of 0 when using add_slide; so we are fine.

prs.save(str(OUT_FILE))
print(f'Written: {OUT_FILE}')
