#!/usr/bin/env python3
"""
Insert generated figures into the existing PPTX and update speaker notes with numeric summary.

Modifies: vibeCoding101/PartX_simulation/presentation/Part3_with_simulation.pptx (saves new file Part3_with_simulation_updated.pptx)

Adds two slides after slide 3 with histogram and ECDF/boxplot.
"""
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
import json

BASE = Path(__file__).resolve().parent
PPTX = BASE / 'Part3_with_simulation.pptx'
OUT = BASE / 'Part3_with_simulation_updated.pptx'
FIG = BASE / 'figures'
SUMMARY = BASE / 'travel_time_comparison' / 'travel_time_summary.json'

if not PPTX.exists():
    raise SystemExit('PPTX not found at ' + str(PPTX))
prs = Presentation(str(PPTX))
# read numeric summary
if SUMMARY.exists():
    with open(SUMMARY,'r') as f:
        summary = json.load(f).get('summary',{})
else:
    summary = {}

# target insertion index: after slide 2 or 3 (0-based). We'll insert after third slide (index 2), append slides then reorder not necessary.
insert_index = 3

# Helper to add image slide
def add_image_slide(title, image_path, notes_text):
    slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    # set title
    if slide.shapes.title:
        slide.shapes.title.text = title
    # add image
    left = Inches(0.5)
    top = Inches(1.4)
    slide.shapes.add_picture(str(image_path), left, top, width=Inches(9))
    slide.notes_slide.notes_text_frame.text = notes_text
    return slide

# create slides
hist = FIG / 'travel_hist_overlay.png'
box = FIG / 'travel_boxplot.png'
cdf = FIG / 'travel_cdf.png'

notes_hist = 'Histogram comparing peak and off-peak travel time distributions. See travel_time_summary.json for numeric summary.'
notes_box = 'Boxplot of travel times. See travel_time_summary.json for numeric summary.'
notes_cdf = 'ECDF of travel times to compare distributional differences between peak and off-peak.'

# Add histogram slide
if hist.exists():
    add_image_slide('Peak vs Off-peak: Travel Time Distribution', hist, notes_hist)
if box.exists():
    add_image_slide('Travel time: Off-peak vs Peak (boxplot)', box, notes_box)
if cdf.exists():
    add_image_slide('ECDF: Peak vs Off-peak', cdf, notes_cdf)

# Add a numeric summary slide
summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
if summary_slide.shapes.title:
    summary_slide.shapes.title.text = 'Summary: travel time comparison (peak vs off-peak)'

body = summary_slide.shapes.placeholders[1].text_frame
# build summary text
speak = []
if 'offpeak' in summary and summary['offpeak']:
    off = summary['offpeak']
    speak.append(f"Off-peak: n={off['count']}, mean={off['mean_s']:.1f}s, median={off['median_s']:.0f}s, p90={off['p90_s']:.0f}s")
if 'peak' in summary and summary['peak']:
    p = summary['peak']
    speak.append(f"Peak: n={p['count']}, mean={p['mean_s']:.1f}s, median={p['median_s']:.0f}s, p90={p['p90_s']:.0f}s")
if 'peak' in summary and 'offpeak' in summary and summary['peak'] and summary['offpeak']:
    diff_mean = summary['peak']['mean_s'] - summary['offpeak']['mean_s']
    diff_med = summary['peak']['median_s'] - summary['offpeak']['median_s']
    diff_p90 = summary['peak']['p90_s'] - summary['offpeak']['p90_s']
    speak.append(f"Differences (peak - offpeak): mean +{diff_mean:.1f}s, median +{diff_med:.0f}s, p90 +{diff_p90:.0f}s")

body.text = '\n'.join(speak)
summary_slide.notes_slide.notes_text_frame.text = 'Numeric summary inserted. Use this slide to report exact numbers.'

prs.save(str(OUT))
print('Saved updated PPTX to', OUT)
